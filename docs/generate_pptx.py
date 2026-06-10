"""
LifeLong WM AI Agent — 발표 자료 자동 생성
실행: conda run -n jbfinai python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
import copy

# ── 색상 팔레트 ───────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1B, 0x4F, 0x8A)
BLUE_MID   = RGBColor(0x4A, 0x9E, 0xDB)
BLUE_LIGHT = RGBColor(0xEB, 0xF2, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x44, 0x44, 0x44)
RED        = RGBColor(0xE5, 0x39, 0x35)
GREEN      = RGBColor(0x2E, 0x7D, 0x32)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 완전 빈 레이아웃


# ── 헬퍼 함수 ────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or GRAY
    return tb

def header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.2, fill=BLUE_DARK)
    add_text(slide, title, 0.4, 0.15, 10, 0.6,
             size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.72, 10, 0.4,
                 size=14, color=BLUE_LIGHT)

def slide_number(slide, n, total=8):
    add_text(slide, f"{n} / {total}", 12.5, 7.1, 0.8, 0.3,
             size=11, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.RIGHT)


# ── Slide 1: 표지 ────────────────────────────────────────────────────────────

s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
add_rect(s1, 0, 4.5, 13.33, 3.0, fill=RGBColor(0x13, 0x3A, 0x6B))

add_text(s1, "LifeLong WM AI Agent", 0.8, 1.0, 11.7, 1.2,
         size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "통계가 불확실성을 계산하고, AI가 누구나 이해할 수 있는 말로 전달한다",
         0.8, 2.2, 11.7, 0.7, size=18, color=BLUE_LIGHT, align=PP_ALIGN.CENTER, italic=True)
add_rect(s1, 4.5, 3.1, 4.33, 0.05, fill=BLUE_MID)

add_text(s1, "JB금융그룹 Fin:AI Challenge  |  지정주제1 개인 라이프케어 AI",
         0.8, 3.4, 11.7, 0.5, size=14, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)
add_text(s1, "2026.06.12", 0.8, 5.2, 11.7, 0.5,
         size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s1, "통계학과 팀  |  Role A: 통계 모델  ·  Role B: AI Agent + UI",
         0.8, 5.8, 11.7, 0.5, size=13, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)


# ── Slide 2: 왜 이 서비스인가 ────────────────────────────────────────────────

s2 = prs.slides.add_slide(BLANK)
header_bar(s2, "왜 이 서비스인가", "기존 WM과의 차별점")
slide_number(s2, 2)

# 왼쪽 박스 — 기존
add_rect(s2, 0.4, 1.4, 5.8, 4.8, fill=RGBColor(0xF5, 0xF5, 0xF5), line=RGBColor(0xCC, 0xCC, 0xCC))
add_text(s2, "기존 WM 서비스", 0.6, 1.5, 5.4, 0.5, size=16, bold=True, color=GRAY)
add_rect(s2, 0.6, 2.0, 5.4, 0.05, fill=RGBColor(0xCC, 0xCC, 0xCC))
for y, txt in zip([2.2, 2.9, 3.6, 4.3],
                  ["❌  미래를 단일 숫자로 표현",
                   "\"83세까지 버팁니다\"",
                   "❌  LLM + 금융 데이터 → 텍스트 조언",
                   "❌  불확실성 무시"]):
    add_text(s2, txt, 0.7, y, 5.2, 0.6, size=15, color=GRAY)

# 오른쪽 박스 — 이 서비스
add_rect(s2, 6.9, 1.4, 6.0, 4.8, fill=BLUE_LIGHT, line=BLUE_DARK)
add_text(s2, "LifeLong WM", 7.1, 1.5, 5.6, 0.5, size=16, bold=True, color=BLUE_DARK)
add_rect(s2, 7.1, 2.0, 5.6, 0.05, fill=BLUE_DARK)
for y, txt in zip([2.2, 2.9, 3.6, 4.3],
                  ["✅  미래를 확률 분포로 표현",
                   "\"83세 전 고갈 확률 41%\"",
                   "✅  통계 모델 → 확률 인사이트 → LLM 해석",
                   "✅  한국 공개 데이터 100% 기반"]):
    add_text(s2, txt, 7.1, y, 5.7, 0.6, size=15, color=BLUE_DARK)

add_text(s2, "→", 6.2, 3.5, 0.6, 0.6, size=30, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)


# ── Slide 3: 전체 아키텍처 ───────────────────────────────────────────────────

s3 = prs.slides.add_slide(BLANK)
header_bar(s3, "전체 아키텍처", "통계 모델 → AI Agent → UI 3계층 구조")
slide_number(s3, 3)

layers = [
    (BLUE_LIGHT,              BLUE_DARK,  "사용자 (Streamlit UI)",
     "버튼 4개  ·  팬 차트  ·  연령별 글씨 크기  ·  음성 모드"),
    (RGBColor(0xE8,0xF5,0xE9), GREEN,     "AI Agent 계층",
     "llama3.1:8b → Tool 선택  ·  exaone3.5:7.8b → 한국어 응답  (완전 로컬, 무료)"),
    (RGBColor(0xFFF3,0xE0,0x00)[:3] if False else RGBColor(0xFF,0xF3,0xE0),
     RGBColor(0xE6,0x5C,0x00),             "통계 모델 계층 (핵심 차별화)",
     "생존분석 · Monte Carlo · 연금최적화 · CVaR · K-means · 이상탐지"),
    (RGBColor(0xF3,0xE5,0xF5), RGBColor(0x6A,0x1B,0x9A), "공개 데이터 계층",
     "통계청 생명표 · 한국은행 ECOS · KRX · 가계금융복지조사 · 국민연금공단"),
]
for i, (bg, fc, title, desc) in enumerate(layers):
    y = 1.35 + i * 1.42
    add_rect(s3, 0.4, y, 12.5, 1.25, fill=bg, line=fc)
    add_text(s3, title, 0.6, y+0.08, 4.5, 0.45, size=15, bold=True, color=fc)
    add_text(s3, desc,  0.6, y+0.55, 12.0, 0.5,  size=13, color=GRAY)
    if i < 3:
        add_text(s3, "▼", 6.2, y+1.22, 1.0, 0.22, size=14, color=BLUE_MID, align=PP_ALIGN.CENTER)


# ── Slide 4: 핵심 기능 ───────────────────────────────────────────────────────

s4 = prs.slides.add_slide(BLANK)
header_bar(s4, "핵심 기능 — 버튼 4개", "어르신이 버튼만 누르면 통계 모델이 자동 실행")
slide_number(s4, 4)

features = [
    ("💰", "내 자산\n얼마나 버티나요?",
     "Monte Carlo 10,000회\n→ 자산 고갈 확률\n→ 팬 차트 시각화"),
    ("📅", "연금 언제 받는 게\n좋을까요?",
     "62~70세 시나리오 비교\n→ 최적 수령 시기\n→ 막대 차트 시각화"),
    ("👥", "또래랑\n비교해주세요",
     "K-means 군집 분석\n→ 퍼센타일 랭킹\n→ 상위 집단 행동 패턴"),
    ("📉", "지출 줄이면\n어떻게 되나요?",
     "지출 20% 절감 재시뮬\n→ 자산 수명 연장 효과\n→ 팬 차트 비교"),
]

for i, (icon, title, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.4
    y = 1.45 + row * 2.8
    add_rect(s4, x, y, 6.0, 2.5, fill=BLUE_LIGHT, line=BLUE_DARK)
    add_text(s4, icon,  x+0.2, y+0.1,  0.8, 0.8,  size=28, align=PP_ALIGN.CENTER)
    add_text(s4, title, x+1.1, y+0.12, 4.7, 0.8,  size=15, bold=True, color=BLUE_DARK)
    add_text(s4, desc,  x+0.2, y+1.0,  5.7, 1.3,  size=13, color=GRAY)


# ── Slide 5: 통계 모델 상세 ──────────────────────────────────────────────────

s5 = prs.slides.add_slide(BLANK)
header_bar(s5, "통계 모델 — 핵심 차별화", "모든 수치는 한국 공개 데이터 기반 확률 계산")
slide_number(s5, 5)

models_info = [
    ("생존 분석", "Cox PH Model", "통계청 국민생명표", "개인화 생존 함수 S(t)"),
    ("Monte Carlo", "10,000회 시뮬레이션", "ECOS 금리/CPI · KRX 수익률", "자산 고갈 확률 · 10/50/90 경로"),
    ("연금 최적화", "시나리오 비교", "국민연금공단 통계", "최적 수령 시기 · 손익분기 연령"),
    ("CVaR 최적화", "cvxpy 볼록 최적화", "KRX 자산 수익률", "리스크 허용도별 자산 배분"),
    ("K-means 군집", "군집 내 퍼센타일", "가계금융복지조사", "또래 대비 재무 위치"),
    ("이상치 탐지", "Isolation Forest · CUSUM", "개인 거래 기저선", "보이스피싱 · 소비 이탈 감지"),
]

headers = ["모델", "방법론", "데이터", "출력"]
col_w = [2.2, 2.8, 3.5, 3.5]
col_x = [0.3, 2.55, 5.4, 8.95]

# 헤더 행
for j, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    add_rect(s5, cx, 1.35, cw-0.05, 0.42, fill=BLUE_DARK)
    add_text(s5, hdr, cx+0.05, 1.38, cw-0.1, 0.35,
             size=13, bold=True, color=WHITE)

for i, (m, method, data, output) in enumerate(models_info):
    y = 1.8 + i * 0.82
    bg = BLUE_LIGHT if i % 2 == 0 else WHITE
    for j, (txt, cx, cw) in enumerate(zip([m, method, data, output], col_x, col_w)):
        add_rect(s5, cx, y, cw-0.05, 0.75, fill=bg, line=RGBColor(0xCC,0xCC,0xCC))
        add_text(s5, txt, cx+0.05, y+0.08, cw-0.1, 0.6,
                 size=12, bold=(j==0), color=BLUE_DARK if j==0 else GRAY)


# ── Slide 6: AI Agent 구조 ───────────────────────────────────────────────────

s6 = prs.slides.add_slide(BLANK)
header_bar(s6, "AI Agent — 이중 모델 구조", "비용 없는 로컬 LLM으로 자연스러운 한국어 응답 구현")
slide_number(s6, 6)

boxes = [
    (0.4,  2.5, 2.5, 2.2, BLUE_LIGHT,              BLUE_DARK,
     "사용자 질문", "버튼 클릭\n(LLM 직접 노출 없음)"),
    (3.3,  2.5, 3.0, 2.2, RGBColor(0xE8,0xF5,0xE9), GREEN,
     "llama3.1:8b", "Tool 선택 · 호출\n(Meta, 영어 강점)"),
    (6.7,  2.5, 2.8, 2.2, RGBColor(0xFF,0xF3,0xE0), RGBColor(0xE6,0x5C,0x00),
     "통계 모델", "순수 Python\nAPI 비용 없음"),
    (9.9,  2.5, 3.0, 2.2, BLUE_LIGHT,              BLUE_DARK,
     "exaone3.5:7.8b", "한국어 설명 생성\n(LG, 한국어 특화)"),
]
arrows = [(2.9, 3.6), (6.3, 3.6), (9.5, 3.6)]

for x, y, w, h, bg, fc, title, desc in boxes:
    add_rect(s6, x, y, w, h, fill=bg, line=fc)
    add_text(s6, title, x+0.1, y+0.1, w-0.2, 0.5, size=14, bold=True, color=fc)
    add_text(s6, desc,  x+0.1, y+0.7, w-0.2, 1.3, size=12, color=GRAY)

for ax, ay in arrows:
    add_text(s6, "→", ax, ay, 0.4, 0.4, size=22, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

add_text(s6, "✅  완전 무료 (Ollama 로컬 실행)  ·  비용 0원",
         0.4, 5.1, 12.5, 0.5, size=15, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_text(s6, "중국어 오염 없음  ·  자연스러운 어르신 말투  ·  연령별 언어 수준 자동 조정",
         0.4, 5.65, 12.5, 0.45, size=13, color=GRAY, align=PP_ALIGN.CENTER)


# ── Slide 7: 평가 지표 ───────────────────────────────────────────────────────

s7 = prs.slides.add_slide(BLANK)
header_bar(s7, "평가 지표", "정량적 검증 — 45개 테스트 케이스 자동 실행")
slide_number(s7, 7)

metrics = [
    ("Tool 선택 정확도", "95%",  "19 / 20", BLUE_DARK, "올바른 통계 도구 선택"),
    ("다중 턴 맥락 유지", "80%",  "4 / 5",  GREEN,     "이전 파라미터 연속 유지"),
    ("위험 차단 Recall",  "73%",  "11 / 15", RED,      "보이스피싱·고위험 조언 차단"),
]

for i, (label, pct, detail, color, desc) in enumerate(metrics):
    x = 0.4 + i * 4.2
    add_rect(s7, x, 1.4, 3.9, 3.6, fill=WHITE, line=color)
    add_rect(s7, x, 1.4, 3.9, 0.55, fill=color)
    add_text(s7, label, x+0.1, 1.45, 3.7, 0.45, size=14, bold=True, color=WHITE)
    add_text(s7, pct,   x+0.1, 2.05, 3.7, 1.4,  size=54, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s7, detail, x+0.1, 3.5, 3.7, 0.45, size=14, color=GRAY, align=PP_ALIGN.CENTER)
    add_text(s7, desc,   x+0.1, 4.0, 3.7, 0.7,  size=12, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s7, "평가 방법: python evaluation/test_cases.py  (자동 실행, 재현 가능)",
         0.4, 5.3, 12.5, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)

add_text(s7, "Coverage Probability · C-index · Silhouette Score 등 통계 모델 지표 별도 측정",
         0.4, 5.75, 12.5, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER, italic=True)


# ── Slide 8: 사회적 가치 & 마무리 ───────────────────────────────────────────

s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
add_rect(s8, 0, 5.8, 13.33, 1.7, fill=RGBColor(0x13, 0x3A, 0x6B))

add_text(s8, "사회적 가치 & 마무리", 0.6, 0.25, 12.0, 0.7,
         size=26, bold=True, color=WHITE)
add_rect(s8, 0.6, 1.0, 5.0, 0.05, fill=BLUE_MID)

points = [
    ("📊", "확률론적 재무 설계",   "단일 숫자가 아닌 확률 분포로 불확실성 정량화"),
    ("🇰🇷", "한국 데이터 100%",    "통계청·한국은행·KRX 공개 데이터 기반"),
    ("👴", "포용적 설계",          "가족 없는 독거 고령층도 완결되는 서비스"),
    ("🔒", "사기 방어 내장",       "보이스피싱·이상 거래 실시간 감지 알림"),
    ("💸", "비용 제로",            "완전 로컬 LLM — Ollama 사용, API 비용 없음"),
]

for i, (icon, title, desc) in enumerate(points):
    y = 1.2 + i * 0.9
    add_text(s8, icon,  0.4, y, 0.6, 0.7, size=20, align=PP_ALIGN.CENTER)
    add_text(s8, title, 1.1, y,      3.5, 0.4, size=15, bold=True,  color=WHITE)
    add_text(s8, desc,  1.1, y+0.38, 11.5, 0.4, size=13, color=BLUE_LIGHT)

add_text(s8, "통계학 전공 역량 × AI Agent = 누구나 이해할 수 있는 확률 기반 재무 동반자",
         0.6, 6.0, 12.1, 0.6, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER, italic=True)


# ── 저장 ─────────────────────────────────────────────────────────────────────

out = "LifeLong_WM_Presentation.pptx"
prs.save(out)
print(f"✅  저장 완료: {out}  ({prs.slides.__len__()} 슬라이드)")
